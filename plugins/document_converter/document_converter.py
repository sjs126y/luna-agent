"""Read-only conversion of common document formats to text or Markdown."""

from __future__ import annotations

import asyncio
import csv
from dataclasses import dataclass, field
from html.parser import HTMLParser
import io
import json
from pathlib import Path
import re
import shutil
import subprocess
import tempfile
from typing import Any

from luna_agent_plugin_sdk import (
    ResourceRequirement,
    ToolResourceBinding,
    ToolEntry,
    ToolHandlerOutput,
)


_DIRECT_SUFFIXES = {".csv", ".docx", ".html", ".htm", ".md", ".pdf", ".pptx", ".txt", ".xlsx"}
_LEGACY_TARGETS = {
    ".doc": ".docx",
    ".odt": ".docx",
    ".rtf": ".docx",
    ".xls": ".xlsx",
    ".ods": ".xlsx",
    ".ppt": ".pptx",
    ".odp": ".pptx",
}
_MAX_FILE_BYTES = 50 * 1024 * 1024
_MAX_CONVERTED_CHARS = 5_000_000
_MAX_PDF_PAGES = 500
_MAX_SLIDES = 500
_MAX_SHEETS = 100
_MAX_ROWS_PER_SHEET = 10_000
_MAX_COLUMNS = 100


class DocumentConvertError(RuntimeError):
    def __init__(self, reason: str, detail: str) -> None:
        self.reason = reason
        self.detail = detail
        super().__init__(detail)


@dataclass
class Conversion:
    content: str
    source_format: str
    details: dict[str, Any] = field(default_factory=dict)
    warnings: list[str] = field(default_factory=list)


def _source_path(raw: str) -> Path:
    requested = Path(str(raw or "")).expanduser()
    if not requested.is_absolute():
        requested = Path.cwd() / requested
    if requested.is_symlink():
        raise DocumentConvertError("document_symlink_blocked", "Symbolic links are not supported.")
    return requested.resolve()


def _resources(input_: dict[str, Any]) -> list[ResourceRequirement]:
    raw = str(input_.get("path") or "").strip()
    if not raw:
        return []
    return [ResourceRequirement("filesystem", str(_source_path(raw)), "read", "document source")]


def _precheck(input_: dict[str, Any]) -> str | None:
    raw = str(input_.get("path") or "").strip()
    if not raw:
        return "A document path is required."
    suffix = Path(raw).suffix.lower()
    if suffix not in _DIRECT_SUFFIXES and suffix not in _LEGACY_TARGETS:
        return f"Unsupported document type: {suffix or 'unknown'}"
    return None


async def document_convert(
    path: str,
    format: str = "markdown",
    offset: int = 0,
    limit: int = 20_000,
) -> str | ToolHandlerOutput:
    try:
        output_format = str(format or "markdown").strip().lower()
        if output_format not in {"markdown", "text"}:
            raise DocumentConvertError("document_format_invalid", "format must be markdown or text")
        source = _source_path(path)
        if not source.exists():
            raise DocumentConvertError("document_missing", f"Document not found: {source}")
        if not source.is_file():
            raise DocumentConvertError("document_not_file", f"Path is not a regular file: {source}")
        size = source.stat().st_size
        if size <= 0:
            raise DocumentConvertError("document_empty", "Document is empty.")
        if size > _MAX_FILE_BYTES:
            raise DocumentConvertError(
                "document_too_large",
                f"Document exceeds the {_MAX_FILE_BYTES} byte limit.",
            )
        converted = await asyncio.to_thread(_convert, source, output_format)
        normalized = _normalize(converted.content)
        if not normalized:
            raise DocumentConvertError("document_no_text", "Document contains no readable text.")
        source_truncated = len(normalized) > _MAX_CONVERTED_CHARS
        if source_truncated:
            normalized = normalized[:_MAX_CONVERTED_CHARS]
            converted.warnings.append(
                f"Converted text was limited to {_MAX_CONVERTED_CHARS} characters."
            )
        start = max(0, int(offset or 0))
        window = max(1, min(int(limit or 20_000), 30_000))
        content = normalized[start : start + window]
        next_offset = start + len(content)
        return json.dumps({
            "ok": True,
            "path": str(source),
            "source_format": converted.source_format,
            "format": output_format,
            "file_size_bytes": size,
            "total_chars": len(normalized),
            "offset": start,
            "returned_chars": len(content),
            "has_more": next_offset < len(normalized),
            "next_offset": next_offset if next_offset < len(normalized) else None,
            "source_truncated": source_truncated,
            "details": converted.details,
            "warnings": converted.warnings,
            "content": content,
        }, ensure_ascii=False, sort_keys=True)
    except DocumentConvertError as exc:
        return ToolHandlerOutput(
            text=json.dumps({
                "ok": False,
                "reason_code": exc.reason,
                "error": exc.detail,
            }, ensure_ascii=False, sort_keys=True),
            metadata={"reason_code": exc.reason},
            is_error=True,
        )
    except Exception as exc:
        return ToolHandlerOutput(
            text=json.dumps({
                "ok": False,
                "reason_code": "document_convert_failed",
                "error": f"{type(exc).__name__}: {exc}",
            }, ensure_ascii=False, sort_keys=True),
            metadata={"reason_code": "document_convert_failed"},
            is_error=True,
        )


def _convert(source: Path, output_format: str) -> Conversion:
    suffix = source.suffix.lower()
    if suffix in _LEGACY_TARGETS:
        return _convert_legacy(source, output_format, _LEGACY_TARGETS[suffix])
    if suffix == ".pdf":
        return _pdf(source, output_format)
    if suffix == ".docx":
        return _docx(source, output_format)
    if suffix == ".pptx":
        return _pptx(source, output_format)
    if suffix == ".xlsx":
        return _xlsx(source, output_format)
    if suffix == ".csv":
        return _csv(source, output_format)
    if suffix in {".html", ".htm"}:
        return _html(source, output_format)
    if suffix in {".txt", ".md"}:
        text = _decode(source.read_bytes())
        if suffix == ".md" and output_format == "text":
            text = _markdown_to_text(text)
        return Conversion(text, suffix.removeprefix("."))
    raise DocumentConvertError("document_type_unsupported", f"Unsupported document type: {suffix}")


def _pdf(source: Path, output_format: str) -> Conversion:
    import fitz

    try:
        document = fitz.open(str(source))
    except Exception as exc:
        raise DocumentConvertError("document_pdf_invalid", str(exc)) from exc
    try:
        if document.needs_pass:
            raise DocumentConvertError("document_password_required", "Password-protected PDF is not supported.")
        pages_read = min(len(document), _MAX_PDF_PAGES)
        sections = []
        for index in range(pages_read):
            text = document[index].get_text("text").strip()
            if not text:
                continue
            header = f"## Page {index + 1}\n\n" if output_format == "markdown" else f"Page {index + 1}\n"
            sections.append(header + text)
        warnings = []
        if pages_read < len(document):
            warnings.append(f"PDF was limited to the first {_MAX_PDF_PAGES} pages.")
        return Conversion(
            "\n\n".join(sections),
            "pdf",
            {"pages_total": len(document), "pages_read": pages_read},
            warnings,
        )
    finally:
        document.close()


def _docx(source: Path, output_format: str) -> Conversion:
    from docx import Document

    try:
        document = Document(str(source))
    except Exception as exc:
        raise DocumentConvertError("document_docx_invalid", str(exc)) from exc
    blocks: list[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style = str(getattr(paragraph.style, "name", "") or "")
        heading = re.fullmatch(r"Heading\s+([1-6])", style, flags=re.IGNORECASE)
        if output_format == "markdown" and heading:
            text = f"{'#' * int(heading.group(1))} {text}"
        elif output_format == "markdown" and style.lower().startswith("list bullet"):
            text = f"- {text}"
        elif output_format == "markdown" and style.lower().startswith("list number"):
            text = f"1. {text}"
        blocks.append(text)
    for table in document.tables:
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        rendered = _table(rows, output_format)
        if rendered:
            blocks.append(rendered)
    return Conversion(
        "\n\n".join(blocks),
        "docx",
        {"paragraphs": len(document.paragraphs), "tables": len(document.tables)},
    )


def _pptx(source: Path, output_format: str) -> Conversion:
    from pptx import Presentation

    try:
        presentation = Presentation(str(source))
    except Exception as exc:
        raise DocumentConvertError("document_pptx_invalid", str(exc)) from exc
    slides_read = min(len(presentation.slides), _MAX_SLIDES)
    sections: list[str] = []
    for index, slide in enumerate(presentation.slides):
        if index >= slides_read:
            break
        items: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                rendered = _table(rows, output_format)
                if rendered:
                    items.append(rendered)
                continue
            text = str(getattr(shape, "text", "") or "").strip()
            if text:
                items.append(text)
        if items:
            header = f"## Slide {index + 1}" if output_format == "markdown" else f"Slide {index + 1}"
            sections.append(header + "\n\n" + "\n\n".join(items))
    warnings = []
    if slides_read < len(presentation.slides):
        warnings.append(f"Presentation was limited to the first {_MAX_SLIDES} slides.")
    return Conversion(
        "\n\n".join(sections),
        "pptx",
        {"slides_total": len(presentation.slides), "slides_read": slides_read},
        warnings,
    )


def _xlsx(source: Path, output_format: str) -> Conversion:
    from openpyxl import load_workbook

    try:
        workbook = load_workbook(source, read_only=True, data_only=True)
    except Exception as exc:
        raise DocumentConvertError("document_xlsx_invalid", str(exc)) from exc
    try:
        sections: list[str] = []
        sheet_names = workbook.sheetnames[:_MAX_SHEETS]
        warnings: list[str] = []
        for name in sheet_names:
            sheet = workbook[name]
            rows = []
            for index, row in enumerate(sheet.iter_rows(values_only=True)):
                if index >= _MAX_ROWS_PER_SHEET:
                    warnings.append(f"Sheet {name} was limited to {_MAX_ROWS_PER_SHEET} rows.")
                    break
                values = [_cell(value) for value in row[:_MAX_COLUMNS]]
                if any(values):
                    rows.append(values)
            if rows:
                header = f"## Sheet: {name}" if output_format == "markdown" else f"Sheet: {name}"
                sections.append(header + "\n\n" + _table(rows, output_format))
        if len(workbook.sheetnames) > _MAX_SHEETS:
            warnings.append(f"Workbook was limited to the first {_MAX_SHEETS} sheets.")
        return Conversion(
            "\n\n".join(sections),
            "xlsx",
            {"sheets_total": len(workbook.sheetnames), "sheets_read": len(sheet_names)},
            warnings,
        )
    finally:
        workbook.close()


def _csv(source: Path, output_format: str) -> Conversion:
    text = _decode(source.read_bytes())
    rows = []
    for index, row in enumerate(csv.reader(io.StringIO(text))):
        if index >= _MAX_ROWS_PER_SHEET:
            break
        rows.append([_cell(value) for value in row[:_MAX_COLUMNS]])
    warnings = [f"CSV was limited to {_MAX_ROWS_PER_SHEET} rows."] if len(rows) >= _MAX_ROWS_PER_SHEET else []
    return Conversion(_table(rows, output_format), "csv", {"rows_read": len(rows)}, warnings)


def _html(source: Path, output_format: str) -> Conversion:
    text = _decode(source.read_bytes())
    if output_format == "markdown":
        import html2text

        converter = html2text.HTML2Text()
        converter.body_width = 0
        converter.ignore_images = True
        converter.ignore_links = False
        content = converter.handle(text)
    else:
        parser = _PlainHTML()
        parser.feed(text)
        content = parser.text()
    return Conversion(content, "html")


def _convert_legacy(source: Path, output_format: str, target_suffix: str) -> Conversion:
    executable = shutil.which("libreoffice") or shutil.which("soffice")
    if not executable:
        raise DocumentConvertError(
            "document_libreoffice_required",
            f"{source.suffix.lower()} conversion requires LibreOffice; use a modern document format instead.",
        )
    with tempfile.TemporaryDirectory(prefix="luna-document-") as directory:
        target_format = target_suffix.removeprefix(".")
        process = subprocess.run(
            [executable, "--headless", "--convert-to", target_format, "--outdir", directory, str(source)],
            capture_output=True,
            text=True,
            timeout=60,
            check=False,
        )
        candidates = list(Path(directory).glob(f"*{target_suffix}"))
        if process.returncode or not candidates:
            detail = (process.stderr or process.stdout or "LibreOffice conversion failed").strip()
            raise DocumentConvertError("document_libreoffice_failed", detail[:1000])
        converted = _convert(candidates[0], output_format)
        converted.source_format = source.suffix.lower().removeprefix(".")
        converted.details["converted_via"] = "libreoffice"
        return converted


def _table(rows: list[list[str]], output_format: str) -> str:
    if not rows:
        return ""
    width = min(max(len(row) for row in rows), _MAX_COLUMNS)
    normalized = [(row + [""] * width)[:width] for row in rows]
    if output_format == "text":
        return "\n".join("\t".join(row).rstrip() for row in normalized)
    escaped = [[cell.replace("|", "\\|").replace("\n", " ") for cell in row] for row in normalized]
    header = escaped[0]
    body = escaped[1:]
    lines = ["| " + " | ".join(header) + " |", "| " + " | ".join("---" for _ in header) + " |"]
    lines.extend("| " + " | ".join(row) + " |" for row in body)
    return "\n".join(lines)


def _cell(value: Any) -> str:
    return "" if value is None else str(value).strip()


def _decode(data: bytes) -> str:
    if b"\x00" in data[:4096]:
        raise DocumentConvertError("document_binary_text", "Text document appears to be binary.")
    for encoding in ("utf-8", "utf-8-sig", "gb18030"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _markdown_to_text(text: str) -> str:
    text = re.sub(r"!\[[^]]*]\([^)]*\)", "", text)
    text = re.sub(r"\[([^]]+)]\([^)]*\)", r"\1", text)
    text = re.sub(r"^\s{0,3}#{1,6}\s+", "", text, flags=re.MULTILINE)
    text = re.sub(r"^\s*[-*+]\s+", "", text, flags=re.MULTILINE)
    text = re.sub(r"[*_`~]", "", text)
    return text


def _normalize(text: str) -> str:
    normalized = str(text or "").replace("\r\n", "\n").replace("\r", "\n")
    normalized = re.sub(r"[ \t]+\n", "\n", normalized)
    normalized = re.sub(r"\n{4,}", "\n\n\n", normalized)
    return normalized.strip()


class _PlainHTML(HTMLParser):
    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.parts: list[str] = []
        self.skip_depth = 0

    def handle_starttag(self, tag: str, attrs) -> None:
        if tag in {"script", "style"}:
            self.skip_depth += 1
        elif not self.skip_depth and tag in {"br", "p", "div", "li", "tr", "h1", "h2", "h3", "h4", "h5", "h6"}:
            self.parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag in {"script", "style"} and self.skip_depth:
            self.skip_depth -= 1
        elif not self.skip_depth and tag in {"p", "div", "li", "tr", "h1", "h2", "h3", "h4", "h5", "h6"}:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        if not self.skip_depth and data.strip():
            self.parts.append(data.strip() + " ")

    def text(self) -> str:
        return "".join(self.parts)


def register(ctx) -> None:
    ctx.register.tool(ToolEntry(
        name="document_convert",
        description=(
            "Convert a local PDF, DOCX, PPTX, XLSX, HTML, CSV, TXT, or Markdown document "
            "to paged plain text or Markdown so its contents can be read. Legacy DOC/XLS/PPT "
            "formats are supported when LibreOffice is installed."
        ),
        schema={
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Allowed local document path."},
                "format": {
                    "type": "string",
                    "enum": ["text", "markdown"],
                    "default": "markdown",
                },
                "offset": {"type": "integer", "minimum": 0, "default": 0},
                "limit": {"type": "integer", "minimum": 1, "maximum": 30000, "default": 20000},
            },
            "required": ["path"],
            "additionalProperties": False,
        },
        handler=document_convert,
        toolset="document",
        permission_category="read",
        tags=["document", "convert", "read", "pdf", "docx", "pptx", "xlsx", "markdown"],
        risk_level="low",
        approval_mode="auto",
        resource_resolver=_resources,
        resource_bindings=(ToolResourceBinding(
            kind="filesystem",
            argument="path",
            access="read",
            reason="document source",
        ),),
        idempotent=True,
        is_parallel_safe=True,
        timeout_seconds=70,
    ))
