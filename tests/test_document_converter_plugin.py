from __future__ import annotations

import json
from pathlib import Path

import pytest

from luna_agent.config import Settings
from luna_agent.plugins import PluginManager, PluginStatus
from luna_agent.tools.registry import dispatch_tool_search, tool_registry
from luna_agent_plugin_sdk import ResourceRequirement


PLUGIN_ROOT = Path(__file__).resolve().parents[1] / "plugins" / "document_converter"
PLUGIN_KEY = "productivity/document-converter"


@pytest.fixture
def loaded_plugin(tmp_path):
    settings = Settings(
        plugin_worker_isolation=False,
        agent_data_dir=tmp_path / "data",
        plugins_dirs=[PLUGIN_ROOT],
        plugins_enabled=[PLUGIN_KEY],
        mcp_enabled=False,
        memory_external_provider="none",
    )
    manager = PluginManager(
        settings,
        plugin_dirs=[PLUGIN_ROOT],
        state_path=tmp_path / "plugin-state.json",
        include_builtin=False,
    )
    manager.load_enabled()
    plugin = manager.list_plugins()[0]
    try:
        yield manager, plugin, plugin.module
    finally:
        manager.unload_plugin(PLUGIN_KEY)


def _payload(result) -> dict:
    assert isinstance(result, str), getattr(result, "text", result)
    return json.loads(result)


@pytest.mark.asyncio
async def test_document_converter_registers_one_read_only_tool(loaded_plugin, tmp_path):
    _, plugin, module = loaded_plugin
    entry = tool_registry.get("document_convert")

    assert plugin.status is PluginStatus.LOADED
    assert plugin.tools_registered == ["document_convert"]
    assert entry is not None
    assert entry.permission_category == "read"
    assert entry.approval_mode == "auto"
    search = json.loads(await dispatch_tool_search("convert docx pdf document to markdown"))
    assert "document_convert" in {item["name"] for item in search["hits"]}
    requirements = entry.resource_resolver({"path": str(tmp_path / "report.pdf")})
    assert requirements == [ResourceRequirement(
        "filesystem",
        str((tmp_path / "report.pdf").resolve()),
        "read",
        "document source",
    )]
    assert module._precheck({"path": "report.exe"}) == "Unsupported document type: .exe"


@pytest.mark.asyncio
async def test_document_converter_reads_docx_as_markdown_and_pages_output(
    loaded_plugin,
    tmp_path,
):
    from docx import Document

    _, _, module = loaded_plugin
    path = tmp_path / "report.docx"
    document = Document()
    document.add_heading("Quarterly Report", level=1)
    document.add_paragraph("Revenue increased this quarter.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Revenue"
    table.cell(1, 1).text = "42"
    document.save(path)

    first = _payload(await module.document_convert(str(path), limit=35))
    second = _payload(await module.document_convert(
        str(path), offset=first["next_offset"], limit=30000
    ))

    assert first["ok"] is True
    assert first["source_format"] == "docx"
    assert first["has_more"] is True
    assert first["next_offset"] == first["returned_chars"]
    combined = first["content"] + second["content"]
    assert "# Quarterly Report" in combined
    assert "Revenue increased" in combined
    assert "| Metric | Value |" in combined
    assert second["has_more"] is False


@pytest.mark.asyncio
async def test_document_converter_reads_pdf_pptx_and_xlsx(loaded_plugin, tmp_path):
    import fitz
    from openpyxl import Workbook
    from pptx import Presentation

    _, _, module = loaded_plugin

    pdf_path = tmp_path / "sample.pdf"
    pdf = fitz.open()
    page = pdf.new_page()
    page.insert_text((72, 72), "PDF document text")
    pdf.save(pdf_path)
    pdf.close()

    pptx_path = tmp_path / "slides.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Roadmap"
    slide.placeholders[1].text = "Ship the converter"
    presentation.save(pptx_path)

    xlsx_path = tmp_path / "budget.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Budget"
    sheet.append(["Item", "Cost"])
    sheet.append(["Hosting", 100])
    workbook.save(xlsx_path)
    workbook.close()

    pdf_result = _payload(await module.document_convert(str(pdf_path), format="text"))
    pptx_result = _payload(await module.document_convert(str(pptx_path)))
    xlsx_result = _payload(await module.document_convert(str(xlsx_path)))

    assert "PDF document text" in pdf_result["content"]
    assert pdf_result["details"]["pages_read"] == 1
    assert "## Slide 1" in pptx_result["content"]
    assert "Ship the converter" in pptx_result["content"]
    assert "## Sheet: Budget" in xlsx_result["content"]
    assert "| Hosting | 100 |" in xlsx_result["content"]


@pytest.mark.asyncio
async def test_document_converter_reads_html_and_rejects_missing_or_legacy_without_libreoffice(
    loaded_plugin,
    tmp_path,
    monkeypatch,
):
    _, _, module = loaded_plugin
    html = tmp_path / "page.html"
    html.write_text(
        "<html><style>hidden</style><h1>Guide</h1><p>Readable body.</p></html>",
        encoding="utf-8",
    )
    legacy = tmp_path / "old.doc"
    legacy.write_bytes(b"legacy")
    monkeypatch.setattr(module.shutil, "which", lambda _: None)

    html_result = _payload(await module.document_convert(str(html), format="text"))
    missing = await module.document_convert(str(tmp_path / "missing.pdf"))
    legacy_result = await module.document_convert(str(legacy))

    assert "Guide" in html_result["content"]
    assert "Readable body." in html_result["content"]
    assert "hidden" not in html_result["content"]
    assert missing.is_error is True
    assert json.loads(missing.text)["reason_code"] == "document_missing"
    assert legacy_result.is_error is True
    assert json.loads(legacy_result.text)["reason_code"] == "document_libreoffice_required"
